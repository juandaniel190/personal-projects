#!/usr/bin/env python3
"""Generate a PowerPoint presentation mirroring cs_apps_presentation.html"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colors ──
NAVY = RGBColor(0x0A, 0x0D, 0x4B)
GREEN = RGBColor(0x31, 0xD8, 0x91)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT = RGBColor(0x66, 0x66, 0x66)
GRAY_BORDER = RGBColor(0xE5, 0xE5, 0xE5)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
WHITE_70 = RGBColor(0xB3, 0xB3, 0xCC)  # approx white 70% on navy
WHITE_90 = RGBColor(0xE6, 0xE6, 0xF0)

# ── Dimensions ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ──
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_font(run, size=10, bold=False, italic=False, color=BLACK, name='Lato'):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name

def add_text(tf, text, size=10, bold=False, italic=False, color=BLACK, align=PP_ALIGN.LEFT, space_after=Pt(4)):
    p = tf.paragraphs[-1] if tf.paragraphs and tf.paragraphs[-1].text == '' and len(tf.paragraphs) == 1 else tf.add_paragraph()
    p.alignment = align
    p.space_after = space_after
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, italic, color)
    return p

def add_rich_paragraph(tf, parts, align=PP_ALIGN.LEFT, space_after=Pt(4)):
    """parts = list of (text, size, bold, italic, color)"""
    p = tf.paragraphs[-1] if tf.paragraphs and tf.paragraphs[-1].text == '' and len(tf.paragraphs) == 1 else tf.add_paragraph()
    p.alignment = align
    p.space_after = space_after
    p.space_before = Pt(0)
    for text, size, bold, italic, color in parts:
        run = p.add_run()
        run.text = text
        set_font(run, size, bold, italic, color)
    return p

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_table(slide, left, top, width, height, rows, cols):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def style_table_cell(cell, text, size=9, bold=False, color=BLACK, align=PP_ALIGN.LEFT, fill_color=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color=color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    # Margins
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)

def style_header_cell(cell, text, size=9, color=BLACK):
    style_table_cell(cell, text, size=size, bold=True, color=color)

def add_shape_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    # Small corner radius
    shape.adjustments[0] = 0.05
    return shape


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Part 1 Cover (dark navy)
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
set_slide_bg(slide, NAVY)

tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True

add_text(tf, "CASE STUDY — DANIEL AMÉZQUITA", size=8, bold=True, color=GREEN, space_after=Pt(16))
add_text(tf, "Part 1: Investment Decision", size=28, bold=True, color=WHITE, space_after=Pt(4))
add_text(tf, "Hypothesis & Data Exploration", size=14, color=WHITE_70, space_after=Pt(20))

# Divider line (green)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.6), Inches(1), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()

add_text(tf, "", space_after=Pt(24))
add_text(tf, "The Product Director wants to know: should we stop, maintain, or increase investment in CS Apps?", size=11, bold=True, color=WHITE_90, space_after=Pt(14))

questions = [
    "1.  What hypotheses frame the investment decision?",
    "2.  What does the data say — and where is data quality a problem?",
    "3.  What 2–3 insights answer the question?",
    "4.  What additional data would deepen the analysis?",
    "5.  What is the strategic recommendation?",
]
for q in questions:
    add_rich_paragraph(tf, [
        (q[:2], 11, False, False, GREEN),
        (q[2:], 11, False, False, WHITE_70),
    ], space_after=Pt(6))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Exec Summary
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

# Title
tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.6))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Revenue is growing but delivery is broken — increase investment to fix before the renewal cliff", size=18, bold=True, color=BLACK)

# Body
tb = add_textbox(slide, MARGIN, Inches(0.9), Inches(12), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "CS Apps shows strong revenue traction (+37% ACV, +30% accounts in 12 months) but faces a structural disconnect between selling, adoption and delivering value. Only 15% of accounts are currently live in production, 71% have zero certified users, and the health rate has declined from 63% to 54% in 10 months. $13.1M in renewal pipeline (38% of total) sits with unhealthy accounts, with a $16.3M Q4 cliff where 44% of UFR is at risk.", size=10, color=GRAY, space_after=Pt(8))
add_text(tf, 'Recommendation: Increase investment — conditionally. The thesis is not "grow faster" but "fix the adoption chain before the renewal cliff converts into churn."', size=10, bold=True, color=BLACK, space_after=Pt(8))
add_text(tf, "Acquisition is strong (revenue is growing) → but adoption is broken (clients aren't going live or engaging) → which puts retention at risk (renewals are deteriorating).", size=10, italic=True, color=GRAY, space_after=Pt(4))

# Framework label
tb = add_textbox(slide, MARGIN, Inches(2.7), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "ANALYTICAL FRAMEWORK", size=8, bold=True, color=GRAY_LIGHT, space_after=Pt(2))
add_text(tf, "The analysis is structured around three pillars that form a causal chain — each pillar feeds the next:", size=10, color=GRAY)

# Three pillar cards
pillars = [
    ("P1", "Acquisition & Growth", "Revenue & account base", "Is CS Apps generating and growing revenue?", "Strong", "+37% ACV, +30% accounts — demand is real"),
    ("P2", "Adoption & Engagement", "Implementation, certification & usage", "Are clients reaching time-to-value, adopting, and engaging?", "Broken", "85% not live, 71% zero certified users, shallow engagement"),
    ("P3", "Retention & ARR", "Renewal pipeline & health", "Will this revenue renew or churn?", "Deteriorating", "Health declining, $13.1M at-risk UFR, flawed health metric"),
]

card_w = Inches(3.8)
card_h = Inches(3.2)
card_top = Inches(3.5)
for idx, (pid, title, sub, question, status, summary) in enumerate(pillars):
    left = MARGIN + idx * (card_w + Inches(0.2))
    shape = add_shape_rect(slide, left, card_top, card_w, card_h, WHITE, GRAY_BORDER)

    tb = add_textbox(slide, left + Pt(12), card_top + Pt(8), card_w - Pt(24), card_h - Pt(16))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, "PILLAR", size=7, bold=True, color=GRAY_LIGHT, space_after=Pt(2))
    add_text(tf, f"{pid} — {title}", size=12, bold=True, color=BLACK, space_after=Pt(2))
    add_text(tf, sub, size=9, color=GRAY, space_after=Pt(10))
    add_text(tf, "CORE QUESTION", size=7, bold=True, color=GRAY_LIGHT, space_after=Pt(2))
    add_text(tf, question, size=9, color=GRAY, space_after=Pt(10))
    add_text(tf, "STATUS", size=7, bold=True, color=GRAY_LIGHT, space_after=Pt(2))
    add_text(tf, status, size=12, bold=True, color=BLACK, space_after=Pt(10))
    add_text(tf, "SUMMARY", size=7, bold=True, color=GRAY_LIGHT, space_after=Pt(2))
    add_text(tf, summary, size=9, color=GRAY)

# Bottom reading guide
tb = add_textbox(slide, MARGIN, Inches(6.85), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Reading this analysis: ", 9, True, False, BLACK),
    ("If acquisition is strong but adoption is broken, then retention will inevitably deteriorate. The investment decision depends on whether the adoption gap is fixable — and whether there is time to fix it before renewals hit.", 9, False, False, GRAY),
])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: P1 Acquisition & Growth
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P1: +37% ACV and +30% accounts in 12 months — demand is real, the problem is not sales", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.75), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P1 Acquisition & Growth — Is CS Apps generating and growing revenue meaningfully?", size=10, italic=True, color=GRAY)

# Left: H1 table
tb = add_textbox(slide, MARGIN, Inches(1.15), Inches(5.5), Inches(0.3))
tf = tb.text_frame
add_text(tf, "H1 — Revenue & Account Growth", size=12, bold=True, color=BLACK)

tbl = add_table(slide, MARGIN, Inches(1.5), Inches(5.5), Inches(0.9), 4, 4)
headers = ["Metric", "Jan 2023", "Dec 2023", "Change"]
data = [
    ["Total CS Apps ACV", "$11.3M", "$15.5M", "+37%"],
    ["Number of Accounts", "134", "174", "+30%"],
    ["Mean ACV per account", "~$84K", "~$89K", "Stable"],
]
for j, h in enumerate(headers):
    style_header_cell(tbl.cell(0, j), h)
for i, row in enumerate(data):
    for j, val in enumerate(row):
        bold = (j == 3)
        style_table_cell(tbl.cell(i+1, j), val, bold=bold, color=BLACK if bold else GRAY, align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

# P1 text
tb = add_textbox(slide, MARGIN, Inches(2.5), Inches(5.5), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Revenue is growing driven by new account acquisition, not price increases. The median ACV remains stable around $60K, suggesting consistent deal sizes with the growth coming from volume.", size=9, color=GRAY)

# Right column: Vertical table
tb = add_textbox(slide, Inches(6.8), Inches(1.15), Inches(6), Inches(0.3))
tf = tb.text_frame
add_text(tf, "H7 — High-Value Segment Opportunity", size=12, bold=True, color=BLACK)

tbl = add_table(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(1.2), 5, 6)
v_headers = ["Vertical", "Accounts", "Mean ACV", "% Healthy", "Share of base", "Total ACV (Dec)"]
for j, h in enumerate(v_headers):
    style_header_cell(tbl.cell(0, j), h, size=8)
verticals = [
    ("Telco", "10", "$234K", "53%", "5.7%", "$2.7M"),
    ("General Retailer", "46", "$92K", "60%", "26.4%", "$4.4M"),
    ("Energy, Utilities & Resources", "8", "$84K", "53%", "4.6%", "$0.8M"),
    ("Food & Beverages", "14", "$82K", "57%", "8.0%", "$1.1M"),
]
for i, row in enumerate(verticals):
    for j, val in enumerate(row):
        style_table_cell(tbl.cell(i+1, j), val, size=8, bold=(j==2), color=BLACK if j==2 else GRAY, align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

# Geo table
tbl2 = add_table(slide, Inches(6.8), Inches(3.0), Inches(6), Inches(0.9), 4, 6)
g_headers = ["Geo", "Accounts", "Mean ACV", "% Healthy", "Share of base", "Total ACV (Dec)"]
for j, h in enumerate(g_headers):
    style_header_cell(tbl2.cell(0, j), h, size=8)
geos = [
    ("Americas", "42", "$113K", "52%", "24.1%", "$4.8M"),
    ("APJ", "11", "$92K", "62%", "6.3%", "$1.4M"),
    ("EMEA", "121", "$74K", "62%", "69.5%", "$9.2M"),
]
for i, row in enumerate(geos):
    for j, val in enumerate(row):
        style_table_cell(tbl2.cell(i+1, j), val, size=8, bold=(j==2), color=BLACK if j==2 else GRAY, align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

# Insight + verdict
tb = add_textbox(slide, Inches(6.8), Inches(4.1), Inches(6), Inches(0.6))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Telco delivers 2.8x the average ACV but only 53% health. Americas is the highest-ACV geo but the least healthy (52%).", size=9, color=GRAY, space_after=Pt(6))

tb = add_textbox(slide, MARGIN, Inches(6.8), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P1 verdict: Acquisition is strong. CS Apps is selling well and growing. The problem is not demand — it's what happens after the sale.", size=10, bold=True, color=BLACK)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: P2 Adoption & Engagement
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P2: 85% of accounts are not live and 71% have zero certified users — adoption is the broken link", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.75), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P2 Adoption & Engagement — Are clients reaching time-to-value, adopting, and engaging?", size=10, italic=True, color=GRAY)

# Left: Implementation status
tb = add_textbox(slide, MARGIN, Inches(1.15), Inches(6), Inches(0.3))
tf = tb.text_frame
add_text(tf, "H4 — Implementation Bottleneck", size=12, bold=True, color=BLACK)

impl_data = [
    ("NULL / Untracked", 53), ("Not started", 8), ("Started", 14),
    ("Partially implemented", 18), ("Implemented", 60), ("Partially lived", 5), ("Lived", 27),
]
impl_colors = [
    RGBColor(0xE0,0xE0,0xE0), RGBColor(0xBD,0xBD,0xBD), RGBColor(0x9E,0x9E,0x9E),
    RGBColor(0x75,0x75,0x75), RGBColor(0x81,0xC7,0x84), RGBColor(0x4C,0xAF,0x50), GREEN,
]
tbl = add_table(slide, MARGIN, Inches(1.5), Inches(5.8), Inches(2.4), 8, 3)
style_header_cell(tbl.cell(0, 0), "Status", size=8)
style_header_cell(tbl.cell(0, 1), "Count", size=8)
style_header_cell(tbl.cell(0, 2), "% of Total", size=8)
for i, (label, count) in enumerate(impl_data):
    pct = round(count / 185 * 100)
    style_table_cell(tbl.cell(i+1, 0), label, size=8, color=BLACK)
    style_table_cell(tbl.cell(i+1, 1), str(count), size=8, bold=True, align=PP_ALIGN.RIGHT)
    style_table_cell(tbl.cell(i+1, 2), f"{pct}%", size=8, align=PP_ALIGN.RIGHT)

tb = add_textbox(slide, MARGIN, Inches(4.0), Inches(5.8), Inches(0.4))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, 'Only 27 (15%) "Lived"; 60 (32%) at "Implemented"; 53 (29%) untracked. 85% not in production.', size=8, color=GRAY)

# Right: Module engagement
tb = add_textbox(slide, Inches(6.8), Inches(1.15), Inches(6), Inches(0.3))
tf = tb.text_frame
add_text(tf, "H6 — Shallow Module Engagement", size=12, bold=True, color=BLACK)

tb = add_textbox(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(0.4))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "1 in 4 sessions is on the Homepage. Users who reach Zoning, Workspace, or Error Analysis show higher engagement.", size=9, color=GRAY)

mod_data = [
    ("Homepage", "25.7%", "2.94", "Navigation, not value"),
    ("Zoning Analysis", "18.2%", "3.81", "Core analytics"),
    ("Journey Analysis", "11.3%", "2.76", "Core analytics"),
    ("Workspace", "7.4%", "4.03", "High stickiness"),
    ("Session Replay", "6.4%", "3.60", "Core analytics"),
    ("Error Analysis", "0.6%", "4.36", "Highest stickiness"),
]
tbl = add_table(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(1.8), 7, 4)
for j, h in enumerate(["Module", "% of Sessions", "Avg Sessions/User", "Note"]):
    style_header_cell(tbl.cell(0, j), h, size=8)
for i, (mod, pct, avg, note) in enumerate(mod_data):
    style_table_cell(tbl.cell(i+1, 0), mod, size=8, bold=True, color=BLACK)
    style_table_cell(tbl.cell(i+1, 1), pct, size=8, align=PP_ALIGN.RIGHT)
    style_table_cell(tbl.cell(i+1, 2), avg, size=8, bold=True, align=PP_ALIGN.RIGHT)
    style_table_cell(tbl.cell(i+1, 3), note, size=8, color=GRAY)

# Certification gap
tb = add_textbox(slide, Inches(6.8), Inches(4.0), Inches(6), Inches(0.3))
tf = tb.text_frame
add_text(tf, "H2 — Certification & Adoption Gap", size=12, bold=True, color=BLACK)

tb = add_textbox(slide, Inches(6.8), Inches(4.35), Inches(6), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
add_rich_paragraph(tf, [("Avg certified users per account: ", 9, True, False, BLACK), ("0.69 (CS Digital: 18.1) — ", 9, False, False, GRAY), ("26x gap", 9, True, False, BLACK)])
add_rich_paragraph(tf, [("Accounts with 0 certified users: ", 9, True, False, BLACK), ("71.4% (132/185)", 9, False, False, GRAY)])
add_rich_paragraph(tf, [("Avg sessions per user per month: ", 9, True, False, BLACK), ("2.84 — Low", 9, False, False, GRAY)])

# Verdict
tb = add_textbox(slide, MARGIN, Inches(6.8), Inches(12), Inches(0.4))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P2 verdict: Adoption is the broken link. 85% of accounts are not live. Those that are live have almost no certified users and shallow engagement.", size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: P3 Retention & ARR
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P3: $13.1M in renewals sits with unhealthy accounts — retention is deteriorating", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.75), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P3 Retention & ARR — Will the revenue CS Apps has created survive renewal, or is it at risk of churn?", size=10, italic=True, color=GRAY)

tb = add_textbox(slide, MARGIN, Inches(1.15), Inches(8), Inches(0.3))
tf = tb.text_frame
add_text(tf, "H5 — Renewal Pipeline at Risk", size=12, bold=True, color=BLACK)

# UFR table
ufr_data = [
    ("FQ 2022-11-01", "$5.0M", "20", "60%", "57%"),
    ("FQ 2023-02-01", "$2.8M", "16", "62%", "58%"),
    ("FQ 2023-05-01", "$5.1M", "23", "87%", "88%"),
    ("FQ 2023-08-01", "$5.4M", "17", "59%", "62%"),
    ("FQ 2023-11-01", "$16.3M", "33", "58%", "56%"),
]
tbl = add_table(slide, MARGIN, Inches(1.5), Inches(8), Inches(1.6), 6, 5)
for j, h in enumerate(["Fiscal Quarter", "Total UFR", "Renewing Accts", "% Accts Healthy", "% UFR Healthy"]):
    style_header_cell(tbl.cell(0, j), h, size=9)
for i, row in enumerate(ufr_data):
    bold = (i == 4)
    for j, val in enumerate(row):
        style_table_cell(tbl.cell(i+1, j), val, size=9, bold=bold, color=BLACK, align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

# Analysis text
tb = add_textbox(slide, Inches(9.2), Inches(1.5), Inches(3.5), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "$13.1M of total UFR (37.9%) sits with unhealthy accounts. The Q4 concentration ($16.3M = 47% of full-year renewal pipeline) amplifies the risk.", size=10, color=GRAY)

# Verdict
tb = add_textbox(slide, Inches(9.2), Inches(5.5), Inches(3.5), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "P3 verdict: Retention is deteriorating. Health is declining, $13.1M in renewals is at risk, and the health metric itself is unreliable.", size=12, bold=True, color=BLACK)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: Action Plan
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Five actions close the delivery gap before the Q4 renewal cliff", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.75), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Five actions close the delivery gap before Q4 — with three data pipeline fixes that are prerequisites", size=10, italic=True, color=GRAY)

tb = add_textbox(slide, MARGIN, Inches(1.1), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Action Plan by Pillar (OKR)", size=12, bold=True, color=BLACK)

# OKR table
tbl = add_table(slide, MARGIN, Inches(1.4), Inches(12), Inches(1.6), 6, 4)
for j, h in enumerate(["Pillar", "Objective", "KPI", "Initiative"]):
    style_header_cell(tbl.cell(0, j), h)
okr_rows = [
    ("P2 Adoption", "Improve engagement and adoption", '"Lived" rate: 15% → 40%+', "Accelerate implementation"),
    ("", "", "0 → 3+ certified users/account", "Drive certification"),
    ("", "", "Homepage share 26% → 15%", "In-product activation flows"),
    ("P1 Acquisition", "Sustain revenue and account growth", "Benchmark Fashion 73% health", "Segment playbooks (Telco & Americas)"),
    ("P3 Retention", "Improve retention, secure UFR", "CS Apps-specific WAU formula", "Rebuild health metric"),
]
for i, (pillar, obj, kpi, init) in enumerate(okr_rows):
    style_table_cell(tbl.cell(i+1, 0), pillar, bold=True, color=BLACK)
    style_table_cell(tbl.cell(i+1, 1), obj, color=GRAY)
    style_table_cell(tbl.cell(i+1, 2), kpi, color=GRAY)
    style_table_cell(tbl.cell(i+1, 3), init, bold=True, color=BLACK)
# Merge P2 cells
tbl.cell(1, 0).merge(tbl.cell(3, 0))
tbl.cell(1, 1).merge(tbl.cell(3, 1))

# Data Quality (left) and Deepen Analysis (right)
# Left: Data Quality
tb = add_textbox(slide, MARGIN, Inches(3.2), Inches(5.8), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Data Quality Audit", size=12, bold=True, color=BLACK)

dq_items = [
    ("NULL Implementation Status", "28.7% — 82 accounts", "HIGH", "518 rows untracked; systemic pipeline issue"),
    ("NULL Contract Dates", "29% start · 28% end", "MEDIUM", "Blocks cohort analysis and time-to-live calculations"),
    ("AVG_WAU_GLOBAL = 0", "156 rows — 49 accounts", "CRITICAL", "Breaks health metric for ~27% of base"),
]
y_pos = Inches(3.55)
for title, stat, impact, desc in dq_items:
    tb = add_textbox(slide, MARGIN, y_pos, Inches(5.8), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    add_rich_paragraph(tf, [
        (title, 8, True, False, BLACK),
        (f"  [{impact}]", 7, True, False, GRAY),
    ], space_after=Pt(1))
    add_text(tf, stat, size=8, bold=True, color=BLACK, space_after=Pt(1))
    add_text(tf, desc, size=7, color=GRAY_LIGHT, space_after=Pt(2))
    y_pos += Inches(0.65)

# Right: Deepen Analysis
tb = add_textbox(slide, Inches(6.8), Inches(3.2), Inches(6), Inches(0.3))
tf = tb.text_frame
add_text(tf, "What Else Would Deepen This Analysis?", size=12, bold=True, color=BLACK)

deepen_items = [
    ("SALES", "Win/loss data by competitor — are we losing to Glassbox or FullStory?"),
    ("SALES", "Cross-sell attach rate — CS Apps penetration into the CS Digital base (growth ceiling)"),
    ("CS TEAM", "Actual renewal outcomes for early cohorts — does unhealthy predict churn?"),
    ("CS TEAM", "Time-to-go-live by segment — is the implementation bottleneck structural or improving?"),
    ("MARKET", "TAM by vertical — is Telco a large market or a saturated niche?"),
    ("CUSTOMER", "Exit interview verbatims from churned accounts — actual reason for leaving"),
]
tb = add_textbox(slide, Inches(6.8), Inches(3.55), Inches(6), Inches(3.5))
tf = tb.text_frame; tf.word_wrap = True
for cat, item in deepen_items:
    add_rich_paragraph(tf, [
        (f"{cat}  ", 7, True, False, GREEN),
        (item, 8, False, False, GRAY),
    ], space_after=Pt(5))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: Part 2 Divider (dark navy)
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
set_slide_bg(slide, NAVY)

tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True

add_text(tf, "ADDITIONAL QUESTION", size=8, bold=True, color=GREEN, space_after=Pt(16))
add_text(tf, "Part 2: AI-Era Pricing & Health", size=28, bold=True, color=WHITE, space_after=Pt(4))
add_text(tf, "Monetization & Account Health Redefinition", size=14, color=WHITE_70, space_after=Pt(20))

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.6), Inches(1), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = GREEN; line.line.fill.background()

add_text(tf, "", space_after=Pt(24))
add_text(tf, "Sense increases value per session but Contentsquare captures none of it. How do we adapt?", size=11, bold=True, color=WHITE_90, space_after=Pt(14))

for q in [
    "1.  Should Sense Analyst be priced as a flat-fee add-on, success-based tier, or credit-based model?",
    "2.  What analysis determines the right pricing model?",
    "3.  How do we redefine account health when AI reduces human time-in-tool?",
]:
    add_rich_paragraph(tf, [(q[:2], 11, False, False, GREEN), (q[2:], 11, False, False, WHITE_70)], space_after=Pt(6))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: Monetization
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Daily Sense Analyst quotas by tier capture AI value without replacing session pricing", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.75), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "How do we capture the extra value Sense delivers per session?", size=10, italic=True, color=GRAY)

tb = add_textbox(slide, MARGIN, Inches(1.1), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Introduce daily Sense Analyst query quotas by tier ", 11, False, False, BLACK),
    ("(Quota + Tier model)", 11, True, False, BLACK),
    (", with Sense Chat included free across all plans.", 11, False, False, BLACK),
])

tb = add_textbox(slide, MARGIN, Inches(1.65), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "How the Market Is Pricing AI Today", size=12, bold=True, color=BLACK)

# Market table
tbl = add_table(slide, MARGIN, Inches(1.95), Inches(12), Inches(1.3), 5, 4)
for j, h in enumerate(["Model", "How It Works", "Companies", "Verdict for CS"]):
    style_header_cell(tbl.cell(0, j), h)
market = [
    ("Flat-Fee Add-On", "Fixed $/year for unlimited AI", "Gainsight (Insight Agent add-on)", "Simple but misaligned — leaves value on table"),
    ("Usage / Token-Based", "Pay per AI query or token", "PostHog, Anthropic API", "Transparent but volatile for budgets"),
    ("Outcome-Based", "Pay per resolved outcome", "Intercom Fin ($0.99/ticket)", "Hard in analytics"),
    ("Quota + Tier", "Daily/monthly cap per plan", "Anthropic Claude Pro, Salesforce", "Best fit — frequent friction drives upgrades"),
]
for i, row in enumerate(market):
    for j, val in enumerate(row):
        style_table_cell(tbl.cell(i+1, j), val, size=8, bold=(j==0), color=BLACK if j==0 else GRAY)

tb = add_textbox(slide, MARGIN, Inches(3.4), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Key findings", size=12, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(3.7), Inches(12), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "• Analytics competitors (Amplitude, Mixpanel, Heap) currently bundle AI into tiers for free. When AI becomes the primary interface, bundling alone won't capture the value.", size=9, color=GRAY, space_after=Pt(6))
add_rich_paragraph(tf, [
    ("• The Anthropic parallel is instructive: ", 9, True, False, BLACK),
    ("Claude Pro users hit a daily usage cap and must either wait or upgrade. This creates frequent, low-friction upgrade pressure without fully blocking work.", 9, False, False, GRAY),
])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: Pricing Model
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "A 5-query daily cap covers 85% of users and creates natural upgrade friction", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.75), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Daily Sense Analyst caps by plan; Sense Chat included free on all tiers", size=10, italic=True, color=GRAY)

# Pricing tiers table
tiers = [
    ("Free", "€0 / forever", "Up to 200K monthly sessions\nSession Replay & heatmaps\nSense Chat included\nSense Analyst not included"),
    ("Growth ★", "From €39 / month", "Starting at 7K monthly sessions\nSense Chat unlimited\nSense Analyst 5 queries/day*\nCovers typical workflow"),
    ("Pro", "Let's talk", "Starting at 1M monthly sessions\nSense Analyst 25 queries/day\nMulti-session replay summaries\nPrecision filtering"),
    ("Enterprise", "Let's talk", "Custom sessions\nSense Analyst unlimited\nError summaries & data feeds\nUnlimited projects"),
]
card_w = Inches(2.95)
for idx, (tier, price, features) in enumerate(tiers):
    left = MARGIN + idx * (card_w + Inches(0.1))
    shape = add_shape_rect(slide, left, Inches(1.15), card_w, Inches(2.8), WHITE, GRAY_BORDER)

    tb = add_textbox(slide, left + Pt(10), Inches(1.25), card_w - Pt(20), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, tier, size=13, bold=True, color=BLACK, space_after=Pt(2))
    add_text(tf, price, size=11, bold=True, color=GREEN if "Growth" in tier else BLACK, space_after=Pt(8))
    for f in features.split('\n'):
        add_rich_paragraph(tf, [("✓ ", 8, False, False, GREEN), (f, 8, False, False, GRAY)], space_after=Pt(2))

# Why this model wins
tb = add_textbox(slide, MARGIN, Inches(4.1), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Why this model wins", size=12, bold=True, color=BLACK)

wins = [
    ("Captures uncaptured value", "Revenue tied to AI usage, not just session volume."),
    ("Preserves session pricing", "Sessions remain the base layer; Sense quota is additive."),
    ("Simple to explain", "Global daily cap per tier is transparent and easy for procurement."),
    ("Natural upgrade path", "Regular friction drives tier upgrades without blocking work."),
]
for idx, (title, text) in enumerate(wins):
    left = MARGIN + idx * Inches(3.1)
    tb = add_textbox(slide, left, Inches(4.45), Inches(2.9), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    add_rich_paragraph(tf, [
        (f"{idx+1}. {title}", 8, True, False, BLACK),
        (f" — {text}", 8, False, False, GRAY),
    ])

# Why 5/day box
shape = add_shape_rect(slide, MARGIN, Inches(5.2), Inches(12), Inches(1.0), RGBColor(0xF8, 0xFA, 0xFC), GRAY_BORDER)
tb = add_textbox(slide, MARGIN + Pt(10), Inches(5.3), Inches(11.8), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
add_rich_paragraph(tf, [
    ("*Why 5/day? ", 8, True, False, BLACK),
    ("On average, customers have 3 sessions daily. If this is a proxy for AI query demand, 5/day covers most users while creating a clear upgrade incentive.", 8, False, False, GRAY),
], space_after=Pt(6))
add_rich_paragraph(tf, [
    ("How it works: ", 8, True, False, BLACK),
    ('A Growth analyst runs 5 Sense Analyst queries in a day. The 6th is blocked with: "Your team hit today\'s 5-query limit. Upgrade to Pro for 25/day, or cap resets tomorrow." No work is lost.', 8, False, False, GRAY),
])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: Data Behind Pricing
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Usage is concentrated — top 20% drive 65% of sessions, validating tiered quotas", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.85), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "Is usage heavily spread across accounts or concentrated among a few? The data shows concentration: top 20% = 65% of sessions, P90 vs. median = 6x.", size=10, color=GRAY)

# Two columns: discarded models + additional data
tb = add_textbox(slide, MARGIN, Inches(1.3), Inches(5.8), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Why did we discard all other models?", size=10, bold=True, italic=True, color=GRAY)

discarded = [
    ("Flat-Fee Add-On (discarded)", "Usage is too concentrated (top 20% = 65%). Heavy users would be subsidized."),
    ("Success-Based Tier (discarded)", "Cannot prove usage drives outcomes with current data."),
    ("Credit-Based (needs more data)", "Cannot test; Sense query logs not available yet."),
    ("Quota + Tier (confirmed)", "Natural daily breakpoints exist (P50, P75, P90). 5/day creates friction without blocking work."),
]
tb = add_textbox(slide, MARGIN, Inches(1.65), Inches(5.8), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
for title, desc in discarded:
    add_rich_paragraph(tf, [(title, 9, True, False, BLACK), (f" — {desc}", 9, False, False, GRAY)], space_after=Pt(6))

tb = add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Additional data that would sharpen this decision?", size=10, bold=True, italic=True, color=GRAY)

additional = [
    ("Sense query logs (type, cost, timestamp)", "Whether different query types cost enough to justify credits."),
    ("Query → action taken (export, share)", "Whether usage correlates with outcomes; would revive Success-Based."),
    ("Renewal outcomes by usage tier", "Whether heavy users actually renew at higher rates."),
]
tb = add_textbox(slide, Inches(6.8), Inches(1.65), Inches(6), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
for title, desc in additional:
    add_rich_paragraph(tf, [(title, 9, True, False, BLACK), (f" — {desc}", 9, False, False, GRAY)], space_after=Pt(6))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: Health KPI
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()

tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "The current health metric masks adoption failure — a three-pillar formula fixes it", size=18, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(0.75), Inches(12), Inches(0.3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "A new three-pillar formula accounts for AI engagement", size=10, italic=True, color=GRAY)

# Formula section
tb = add_textbox(slide, MARGIN, Inches(1.15), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Formula", size=12, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(1.5), Inches(5.5), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "If Non-Live / Partially Live:", size=9, bold=True, color=BLACK, space_after=Pt(2))
add_text(tf, "Score = (0.15 x P1) + (0.25 x P2) + (0.60 x P3)", size=9, color=BLACK, space_after=Pt(6))

tb = add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "If Live:", size=9, bold=True, color=BLACK, space_after=Pt(2))
add_text(tf, "Score = (0.3 x P2) + (0.7 x P3)", size=9, color=BLACK, space_after=Pt(6))

tb = add_textbox(slide, MARGIN, Inches(2.15), Inches(12), Inches(0.25))
tf = tb.text_frame
add_text(tf, "Healthy = Score >= 50", size=10, bold=True, color=BLACK)

# Pillars
tb = add_textbox(slide, MARGIN, Inches(2.5), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Pillars", size=12, bold=True, color=BLACK)

pillar_defs = [
    ("P1 — Implementation Velocity", "Non-Live only. Dropped once Live.", "100 if within window; linear decay beyond"),
    ("P2 — Adoption (ACV-Normalized)", "min(Cert/(ACV x a), 1) x 100", "Normalizes adoption vs account size"),
    ("P3 — Engagement Quality (AI-Aware)", "max(0, 100 - ACV/(WAU + Sense Q/wk) x b)", "WAU + Sense Queries/week"),
]
for idx, (title, formula, desc) in enumerate(pillar_defs):
    left = MARGIN + idx * Inches(4.1)
    shape = add_shape_rect(slide, left, Inches(2.85), Inches(3.9), Inches(1.1), WHITE, GRAY_BORDER)
    tb = add_textbox(slide, left + Pt(8), Inches(2.9), Inches(3.7), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, title, size=9, bold=True, color=BLACK, space_after=Pt(3))
    add_text(tf, formula, size=8, italic=True, color=GRAY, space_after=Pt(3))
    add_text(tf, desc, size=8, color=GRAY)

# Examples table
tb = add_textbox(slide, MARGIN, Inches(4.1), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Illustrative Examples (Dummy Data)", size=12, bold=True, color=BLACK)

tb = add_textbox(slide, MARGIN, Inches(4.35), Inches(12), Inches(0.2))
tf = tb.text_frame
add_text(tf, "Assumed constants: a = 1/100,000 · b = 1/500 · Expected implementation window = 90 days", size=8, italic=True, color=GRAY)

tbl = add_table(slide, MARGIN, Inches(4.6), Inches(12), Inches(0.7), 3, 11)
ex_headers = ["Account", "ACV", "Status", "Impl Days", "WAU", "Sense Q/wk", "Cert Users", "P1", "P2", "P3", "Score"]
for j, h in enumerate(ex_headers):
    style_header_cell(tbl.cell(0, j), h, size=7)
examples = [
    ("A — Acme Corp", "$200K", "Non-Live", "210", "2", "0", "1", "0", "50", "0", "15 ⚠"),
    ("C — Gamma Inc", "$200K", "Live", "—", "20", "10", "5", "—", "100", "87", "92 ✓"),
]
for i, row in enumerate(examples):
    for j, val in enumerate(row):
        style_table_cell(tbl.cell(i+1, j), val, size=7, bold=(j==10), color=BLACK, align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

# Data Required table
tb = add_textbox(slide, MARGIN, Inches(5.5), Inches(12), Inches(0.3))
tf = tb.text_frame
add_text(tf, "Data Required to Validate", size=12, bold=True, color=BLACK)

tbl = add_table(slide, MARGIN, Inches(5.8), Inches(12), Inches(1.4), 6, 3)
for j, h in enumerate(["Input", "Current Estimate", "Needed To Calibrate"]):
    style_header_cell(tbl.cell(0, j), h, size=8)
validation_data = [
    ("a — adoption rate per $ ACV", "1/100,000", "Distribution of (Cert Users / ACV) across live accounts"),
    ("b — engagement sensitivity", "1/500", "Distribution of ACV / (WAU + Sense Q) ratios"),
    ("Break-even ratio (P3 = 50)", "~$25K ACV per (WAU+Q)", "Recalibrate once Sense query telemetry available"),
    ("Expected implementation window", "TBD", "Historical time-to-go-live by account tier"),
    ("Sense Queries/week", "Not yet available", "Requires Sense query logs at account level"),
]
for i, row in enumerate(validation_data):
    for j, val in enumerate(row):
        style_table_cell(tbl.cell(i+1, j), val, size=7, bold=(j==0), color=BLACK if j==0 else GRAY)


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = "/Users/damezquita/Documents/GitHub/personal-projects/business-cases/Contentsquare/CS_Apps_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
